"""Format-specific internals for the Media specialist.

The LLM sees a small tool surface in :mod:`mounir.specialists.media`.  This
module owns format detection, bounded extraction, deterministic document
creation, and provider adapters.  Optional dependencies are imported lazily so
one missing format library does not disable the specialist.
"""

from __future__ import annotations

import base64
import csv
import io
import json
import mimetypes
import os
import re
import subprocess
import tempfile
import wave
from pathlib import Path
from typing import Callable
from urllib.parse import urlsplit, urlunsplit
from xml.sax.saxutils import escape

import requests

from .. import path_search


IMAGE_EXT = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}
AUDIO_EXT = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac"}
VIDEO_EXT = {".mp4", ".mov", ".mkv", ".webm", ".avi", ".m4v"}
PRESENTATION_EXT = {".pptx"}
PDF_EXT = {".pdf"}
SPREADSHEET_EXT = {".xlsx", ".xlsm", ".csv", ".tsv"}
DOCUMENT_EXT = {".docx"}
TEXT_EXT = {
    ".txt", ".md", ".rst", ".log", ".json", ".jsonl", ".yaml", ".yml",
    ".xml", ".html", ".htm", ".ini", ".cfg", ".toml", ".sql",
    ".py", ".pyi", ".js", ".jsx", ".ts", ".tsx", ".css", ".scss",
    ".sass", ".less", ".sh", ".bash", ".zsh", ".fish", ".ps1", ".bat",
    ".c", ".h", ".cc", ".cpp", ".hpp", ".cs", ".java", ".kt", ".kts",
    ".go", ".rs", ".rb", ".php", ".swift", ".scala", ".lua", ".r",
    ".dockerfile", ".env", ".gitignore", ".gitattributes", ".editorconfig",
}
TEXT_FILENAMES = {
    "dockerfile", "makefile", "procfile", "license", "readme", "agents.md",
}
FILE_EXT = PDF_EXT | SPREADSHEET_EXT | DOCUMENT_EXT | TEXT_EXT
MEDIA_EXT = IMAGE_EXT | AUDIO_EXT | VIDEO_EXT | PRESENTATION_EXT
ALL_EXT = FILE_EXT | MEDIA_EXT

MAX_TEXT_CHARS = 18_000
MAX_SHEETS = 12
MAX_ROWS_PER_SHEET = 150
MAX_COLUMNS_PER_SHEET = 40
MAX_PRESENTATION_SLIDES = 40
MAX_PRESENTATION_IMAGES = 10
PDF_MIN_TEXT_CHARS = 80
PDF_RENDER_MAX_PAGES = 8
VIDEO_DEFAULT_FRAMES = 8
VIDEO_MAX_FRAMES = 12
IMAGE_MAX_SIDE = 1568
IMAGE_MAX_BYTES = 180_000
AUDIO_INLINE_MAX_BYTES = 20 * 1024 * 1024

MediaResult = tuple[str, list[dict]]
_FILES_READ: set[str] = set()


def reset_task_state() -> None:
    """Require fresh reads before edits in each delegated task."""
    _FILES_READ.clear()


def kind(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in IMAGE_EXT:
        return "image"
    if ext in AUDIO_EXT:
        return "audio"
    if ext in VIDEO_EXT:
        return "video"
    if ext in PRESENTATION_EXT:
        return "presentation"
    if ext in PDF_EXT:
        return "pdf"
    if ext in SPREADSHEET_EXT:
        return "spreadsheet"
    if ext in DOCUMENT_EXT:
        return "document"
    if ext in TEXT_EXT or path.name.casefold() in TEXT_FILENAMES:
        return "text"
    return "unknown"


def resolve_existing(path: str) -> tuple[Path | None, str | None]:
    resolution = path_search.resolve_existing(path, "file")
    return resolution.path, resolution.message or None


def _bounded(text: str, limit: int = MAX_TEXT_CHARS) -> str:
    value = text.strip()
    if len(value) <= limit:
        return value
    return value[:limit].rstrip() + "\n… [truncated]"


def _image_part(data: bytes, mime: str) -> dict:
    encoded = base64.b64encode(data).decode("ascii")
    return {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{encoded}"}}


def _shrink_image(data: bytes) -> tuple[bytes, str]:
    try:
        from PIL import Image
    except ImportError:
        return data, "image/jpeg"

    image = Image.open(io.BytesIO(data))
    if image.mode in ("RGBA", "P", "LA"):
        image = image.convert("RGB")
    image.thumbnail((IMAGE_MAX_SIDE, IMAGE_MAX_SIDE))
    quality = 85
    while True:
        buffer = io.BytesIO()
        image.save(buffer, format="JPEG", quality=quality)
        result = buffer.getvalue()
        if len(result) <= IMAGE_MAX_BYTES or quality <= 35:
            return result, "image/jpeg"
        quality -= 15


def _read_text(
    path: Path, start_line: int = 1, end_line: int | None = None
) -> MediaResult:
    try:
        lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    except Exception as exc:
        return f"Could not read {path}: {exc}", []
    _FILES_READ.add(str(path.resolve()))
    if not lines:
        return f"{path} is empty.", []
    start = max(1, int(start_line))
    if start > len(lines):
        return f"{path} has {len(lines)} lines; line {start} is past the end.", []
    end = min(len(lines), start + 299) if end_line is None else min(len(lines), int(end_line))
    body = "\n".join(
        f"{line_number:>5}\t{line}"
        for line_number, line in enumerate(lines[start - 1:end], start=start)
    )
    body = _bounded(body, 12_000)
    more = (
        f"\n… {len(lines) - end} more line(s); read again from line {end + 1}."
        if end < len(lines) else ""
    )
    return f"{path} (lines {start}-{end} of {len(lines)}):\n{body}{more}", []


def _read_delimited(path: Path) -> MediaResult:
    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    try:
        with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
            rows = []
            for index, row in enumerate(csv.reader(handle, delimiter=delimiter)):
                if index >= MAX_ROWS_PER_SHEET:
                    break
                rows.append(row[:MAX_COLUMNS_PER_SHEET])
    except Exception as exc:
        return f"Could not read {path}: {exc}", []
    body = "\n".join("\t".join(str(cell) for cell in row) for row in rows)
    return (
        f"Read {path.name} ({len(rows)} displayed row(s)):\n\n"
        f"{_bounded(body) or '(file is empty)'}",
        [],
    )


def _read_workbook(path: Path) -> MediaResult:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "Can't read Excel workbooks: install 'openpyxl'.", []
    try:
        workbook = load_workbook(
            path, read_only=True, data_only=False, keep_links=False
        )
        sections = []
        for sheet in workbook.worksheets[:MAX_SHEETS]:
            lines = [f"Sheet: {sheet.title}"]
            for row_index, row in enumerate(
                sheet.iter_rows(max_col=MAX_COLUMNS_PER_SHEET), start=1
            ):
                if row_index > MAX_ROWS_PER_SHEET:
                    lines.append("… [more rows omitted]")
                    break
                values = ["" if cell.value is None else str(cell.value) for cell in row]
                while values and not values[-1]:
                    values.pop()
                if values:
                    lines.append("\t".join(values))
            sections.append("\n".join(lines))
        workbook.close()
    except Exception as exc:
        return f"Could not read {path.name}: {exc}", []
    return (
        f"Read workbook {path.name}:\n\n{_bounded(chr(10).join(sections))}",
        [],
    )


def _read_docx(path: Path) -> MediaResult:
    try:
        from docx import Document
    except ImportError:
        return "Can't read Word documents: install 'python-docx'.", []
    try:
        document = Document(path)
        blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table_index, table in enumerate(document.tables, start=1):
            blocks.append(f"Table {table_index}:")
            blocks.extend(
                "\t".join(cell.text.strip() for cell in row.cells)
                for row in table.rows
            )
    except Exception as exc:
        return f"Could not read {path.name}: {exc}", []
    return f"Read document {path.name}:\n\n{_bounded(chr(10).join(blocks))}", []


def _read_pdf(path: Path) -> MediaResult:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "Can't read PDFs: install 'pypdf'.", []
    try:
        reader = PdfReader(str(path))
        text = "\n".join((page.extract_text() or "") for page in reader.pages).strip()
    except Exception as exc:
        return f"Could not open {path.name}: {exc}", []
    if len(text) >= PDF_MIN_TEXT_CHARS:
        return (
            f"Extracted text from {path.name} ({len(reader.pages)} pages):\n\n{_bounded(text)}",
            [],
        )
    try:
        import fitz
    except ImportError:
        return (
            f"{path.name} has no selectable text. Install 'PyMuPDF' to read it as images.",
            [],
        )
    parts = []
    with fitz.open(str(path)) as document:
        total = len(document)
        for page in document[:PDF_RENDER_MAX_PAGES]:
            data, mime = _shrink_image(page.get_pixmap(dpi=150).tobytes("png"))
            parts.append(_image_part(data, mime))
    return (
        f"{path.name} is scanned; attached {len(parts)} of {total} page(s) for analysis.",
        parts,
    )


def _looks_like_text(path: Path) -> bool:
    try:
        sample = path.read_bytes()[:4096]
    except OSError:
        return False
    if b"\x00" in sample:
        return False
    try:
        sample.decode("utf-8")
        return True
    except UnicodeDecodeError:
        mime = mimetypes.guess_type(path.name)[0] or ""
        return mime.startswith("text/")


def read_file(
    path: str, start_line: int = 1, end_line: int | None = None
) -> MediaResult:
    resolved, error = resolve_existing(path)
    if error:
        return error, []
    assert resolved is not None
    file_kind = kind(resolved)
    if file_kind in {"image", "audio", "video", "presentation"}:
        return f"{resolved.name} is {file_kind} media; use load_media instead.", []
    if resolved.suffix.lower() in {".csv", ".tsv"}:
        return _read_delimited(resolved)
    if resolved.suffix.lower() in {".xlsx", ".xlsm"}:
        return _read_workbook(resolved)
    if file_kind == "document":
        return _read_docx(resolved)
    if file_kind == "pdf":
        return _read_pdf(resolved)
    if file_kind == "text" or _looks_like_text(resolved):
        return _read_text(resolved, start_line, end_line)
    mime = mimetypes.guess_type(resolved.name)[0] or "unknown"
    return f"Unsupported file type {resolved.suffix or '(none)'} ({mime}).", []


def _load_image(path: Path) -> MediaResult:
    try:
        data, mime = _shrink_image(path.read_bytes())
    except Exception as exc:
        return f"Could not load {path.name}: {exc}", []
    if len(base64.b64encode(data)) > IMAGE_MAX_BYTES * 1.4:
        return f"{path.name} is too large to inline; install Pillow to downscale it.", []
    return f"Loaded image {path.name}.", [_image_part(data, mime)]


def _audio_part(path: Path) -> dict:
    audio_format = path.suffix.lower().lstrip(".")
    if audio_format == "m4a":
        audio_format = "mp4"
    return {
        "type": "input_audio",
        "input_audio": {
            "data": base64.b64encode(path.read_bytes()).decode("ascii"),
            "format": audio_format,
        },
    }


def _transcribe(path: Path) -> tuple[str, str]:
    """Use Mounir's saved STT configuration; return text and an error detail."""
    try:
        import numpy as np

        from .. import stt
    except ImportError as exc:
        return "", f"transcription dependency unavailable: {exc}"
    descriptor, wav_path = tempfile.mkstemp(suffix=".wav")
    os.close(descriptor)
    try:
        process = subprocess.run(
            [
                "ffmpeg", "-nostdin", "-loglevel", "error", "-y", "-i", str(path),
                "-ar", "16000", "-ac", "1", "-f", "wav", wav_path,
            ],
            capture_output=True,
            timeout=120,
        )
        if process.returncode:
            detail = process.stderr.decode(errors="replace")[-300:].strip()
            return "", f"ffmpeg could not decode audio: {detail}"
        with wave.open(wav_path, "rb") as handle:
            pcm = np.frombuffer(handle.readframes(handle.getnframes()), dtype=np.int16)
        text, language = stt.transcribe(pcm.astype(np.float32) / 32768.0)
        return text.strip(), f"language={language}" if language else ""
    except FileNotFoundError:
        return "", "ffmpeg is not installed"
    except Exception as exc:
        return "", str(exc)
    finally:
        Path(wav_path).unlink(missing_ok=True)


def _sample_video(path: Path, count: int = VIDEO_DEFAULT_FRAMES) -> MediaResult:
    try:
        import cv2
    except ImportError:
        return "Can't inspect video frames: install 'opencv-python-headless'.", []
    capture = cv2.VideoCapture(str(path))
    total = int(capture.get(cv2.CAP_PROP_FRAME_COUNT)) or 0
    if total <= 0:
        capture.release()
        return f"Could not read frames from {path.name}.", []
    count = max(1, min(int(count), VIDEO_MAX_FRAMES))
    positions = sorted(
        {
            round(index * (total - 1) / max(count - 1, 1))
            for index in range(count)
        }
    )
    parts = []
    for position in positions:
        capture.set(cv2.CAP_PROP_POS_FRAMES, position)
        ok, frame = capture.read()
        if not ok:
            continue
        ok, buffer = cv2.imencode(".jpg", frame)
        if not ok:
            continue
        data, mime = _shrink_image(buffer.tobytes())
        parts.append(_image_part(data, mime))
    capture.release()
    transcript, transcript_note = _transcribe(path)
    summary = f"Loaded video {path.name}; attached {len(parts)} representative frame(s)."
    if transcript:
        summary += f"\n\nAudio transcript:\n{_bounded(transcript)}"
    else:
        summary += f"\n\nAudio transcript unavailable: {transcript_note or 'no speech detected'}."
    return summary, parts


def _load_audio(path: Path) -> MediaResult:
    transcript, note = _transcribe(path)
    if transcript:
        label = f" ({note})" if note else ""
        return f"Transcribed {path.name}{label}:\n\n{_bounded(transcript)}", []
    if path.stat().st_size > AUDIO_INLINE_MAX_BYTES:
        return (
            f"Could not transcribe {path.name}: {note}. The file is too large "
            "to send inline to the selected model.",
            [],
        )
    try:
        part = _audio_part(path)
    except Exception as exc:
        return f"Could not load {path.name}: {exc}", []
    return f"Loaded audio {path.name}; transcription unavailable: {note}.", [part]


def _load_presentation(path: Path) -> MediaResult:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return "Can't read presentations: install 'python-pptx'.", []
    try:
        presentation = Presentation(path)
        sections = []
        parts = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            if slide_index > MAX_PRESENTATION_SLIDES:
                break
            texts = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    value = shape.text.strip()
                    if value:
                        texts.append(value)
                if (
                    shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    and len(parts) < MAX_PRESENTATION_IMAGES
                ):
                    data, mime = _shrink_image(shape.image.blob)
                    parts.append(_image_part(data, mime))
            sections.append(f"Slide {slide_index}:\n" + ("\n".join(texts) or "(no text)"))
    except Exception as exc:
        return f"Could not read {path.name}: {exc}", []
    return (
        f"Read presentation {path.name} ({len(presentation.slides)} slides):\n\n"
        f"{_bounded(chr(10).join(sections))}\n\n"
        f"Attached {len(parts)} embedded image(s) for visual analysis.",
        parts,
    )


def load_media(path: str) -> MediaResult:
    resolved, error = resolve_existing(path)
    if error:
        return error, []
    assert resolved is not None
    media_kind = kind(resolved)
    if media_kind == "image":
        return _load_image(resolved)
    if media_kind == "audio":
        return _load_audio(resolved)
    if media_kind == "video":
        return _sample_video(resolved)
    if media_kind == "presentation":
        return _load_presentation(resolved)
    return f"{resolved.name} is a {media_kind} file; use read_file instead.", []


def find_files(
    directory: str = ".",
    query: str = "",
    group: str = "any",
    recursive: bool = False,
) -> MediaResult:
    resolution = path_search.resolve_existing(directory or ".", "directory")
    if resolution.path is None:
        return resolution.message, []
    folder = resolution.path
    normalized_group = group.lower().strip()
    if normalized_group not in {"file", "media", "directory", "any"}:
        return "group must be file, media, directory, or any.", []
    needle = query.casefold().strip()

    def accepted(item: Path) -> bool:
        if normalized_group == "directory":
            return item.is_dir()
        if normalized_group == "media":
            return item.is_file() and item.suffix.lower() in MEDIA_EXT
        if normalized_group == "file":
            return item.is_file() and item.suffix.lower() not in MEDIA_EXT
        return item.is_file() or item.is_dir()

    try:
        immediate = [
            item for item in folder.iterdir()
            if accepted(item) and (not needle or needle in item.name.casefold())
        ]
    except Exception as exc:
        return f"Could not list {folder}: {exc}", []
    matches = immediate
    searched_recursively = False
    if needle:
        expected = "directory" if normalized_group == "directory" else "any"
        fuzzy_matches = [
            item for item in path_search.find_matches(folder, query, expected)
            if accepted(item)
        ]
        if recursive or not matches:
            matches = fuzzy_matches
        else:
            matches = list(dict.fromkeys([*matches, *fuzzy_matches]))
        searched_recursively = any(item.parent != folder for item in matches)
    matches = sorted(
        matches,
        key=lambda item: (
            not item.is_dir(),
            -item.stat().st_mtime,
            item.name.casefold(),
        ),
    )
    if not matches:
        suffix = f" matching '{query}'" if query else ""
        return f"No {normalized_group} matches{suffix} in {folder}.", []
    scope = "recursively" if searched_recursively else "directly"
    lines = [f"Matches {scope} in {folder}:"]
    for item in matches[:50]:
        item_kind = "directory" if item.is_dir() else kind(item)
        lines.append(f"  {item}  ({item_kind})")
    return "\n".join(lines), []


def _content_object(content: str):
    value = str(content or "")
    try:
        return json.loads(value)
    except (TypeError, ValueError):
        return value


def _write_text(destination: Path, content: str) -> None:
    if destination.suffix.lower() == ".json":
        parsed = json.loads(content)
        content = json.dumps(parsed, ensure_ascii=False, indent=2)
    destination.write_text(content, encoding="utf-8")


def _write_pdf(destination: Path, content: str) -> None:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
    except ImportError as exc:
        raise RuntimeError("install 'reportlab' to create PDFs") from exc
    styles = getSampleStyleSheet()
    story = []
    for block in str(content).split("\n\n"):
        rendered = escape(block.strip()).replace("\n", "<br/>")
        if rendered:
            story.extend((Paragraph(rendered, styles["BodyText"]), Spacer(1, 3 * mm)))
    if not story:
        story.append(Paragraph(" ", styles["BodyText"]))
    SimpleDocTemplate(str(destination), pagesize=A4).build(story)


def _sheet_specs(content: str) -> list[dict]:
    parsed = _content_object(content)
    if isinstance(parsed, list):
        return [{"name": "Sheet1", "rows": parsed}]
    if isinstance(parsed, dict) and isinstance(parsed.get("sheets"), list):
        return parsed["sheets"]
    if isinstance(parsed, dict):
        return [{"name": name, "rows": rows} for name, rows in parsed.items()]
    rows = list(csv.reader(io.StringIO(str(parsed))))
    return [{"name": "Sheet1", "rows": rows}]


def _write_workbook(destination: Path, content: str) -> None:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font
    except ImportError as exc:
        raise RuntimeError("install 'openpyxl' to create Excel workbooks") from exc
    workbook = Workbook()
    workbook.remove(workbook.active)
    for index, spec in enumerate(_sheet_specs(content), start=1):
        if not isinstance(spec, dict):
            raise ValueError(f"sheet {index} must be an object")
        title = str(spec.get("name") or f"Sheet{index}")[:31]
        sheet = workbook.create_sheet(title=title)
        rows = spec.get("rows") or []
        if not isinstance(rows, list):
            raise ValueError(f"rows for {title} must be an array")
        for row in rows:
            sheet.append(row if isinstance(row, list) else [row])
        if rows:
            for cell in sheet[1]:
                cell.font = Font(bold=True)
        for column in sheet.columns:
            width = min(
                max((len(str(cell.value or "")) for cell in column), default=8) + 2,
                60,
            )
            sheet.column_dimensions[column[0].column_letter].width = width
    if not workbook.worksheets:
        workbook.create_sheet("Sheet1")
    workbook.save(destination)


def _write_docx(destination: Path, content: str) -> None:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("install 'python-docx' to create Word documents") from exc
    parsed = _content_object(content)
    document = Document()
    if isinstance(parsed, dict):
        if parsed.get("title"):
            document.add_heading(str(parsed["title"]), level=0)
        for section in parsed.get("sections") or []:
            if isinstance(section, dict):
                if section.get("heading"):
                    document.add_heading(str(section["heading"]), level=1)
                for paragraph in section.get("paragraphs") or []:
                    document.add_paragraph(str(paragraph))
        for table_spec in parsed.get("tables") or []:
            if not isinstance(table_spec, dict):
                continue
            rows = table_spec.get("rows") or []
            headers = table_spec.get("headers") or []
            width = max(len(headers), max((len(row) for row in rows), default=0), 1)
            table = document.add_table(rows=1 if headers else 0, cols=width)
            if headers:
                for index, value in enumerate(headers):
                    table.rows[0].cells[index].text = str(value)
            for row in rows:
                cells = table.add_row().cells
                for index, value in enumerate(row[:width]):
                    cells[index].text = str(value)
    else:
        for block in str(parsed).split("\n\n"):
            document.add_paragraph(block)
    document.save(destination)


def _write_presentation(destination: Path, content: str) -> None:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("install 'python-pptx' to create presentations") from exc
    parsed = _content_object(content)
    if not isinstance(parsed, dict) or not isinstance(parsed.get("slides"), list):
        blocks = [block.strip() for block in str(content).split("\n\n") if block.strip()]
        parsed = {
            "title": blocks[0] if blocks else "Presentation",
            "slides": [
                {"title": f"Slide {index}", "bullets": block.splitlines()}
                for index, block in enumerate(blocks[1:], start=2)
            ],
        }
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = str(parsed.get("title") or "Presentation")
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = str(parsed.get("subtitle") or "")
    for spec in parsed.get("slides") or []:
        if not isinstance(spec, dict):
            continue
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(spec.get("title") or "")
        frame = slide.placeholders[1].text_frame
        frame.clear()
        bullets = spec.get("bullets") or spec.get("content") or []
        if isinstance(bullets, str):
            bullets = bullets.splitlines()
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = str(bullet)
    presentation.save(destination)


def _atomic_create(path: Path, writer: Callable[[Path], None]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    descriptor, temporary = tempfile.mkstemp(
        prefix=f".{path.stem}-", suffix=path.suffix, dir=path.parent
    )
    os.close(descriptor)
    temporary_path = Path(temporary)
    try:
        writer(temporary_path)
        if not temporary_path.is_file() or temporary_path.stat().st_size == 0:
            raise RuntimeError("generator produced an empty file")
        temporary_path.replace(path)
    finally:
        temporary_path.unlink(missing_ok=True)


def create_file(path: str, content: str) -> str:
    from .. import tools

    resolution = path_search.resolve_output(path)
    if resolution.path is None:
        return resolution.message
    destination = resolution.path
    blocked = tools.write_path_block_reason(destination)
    if blocked:
        return blocked
    extension = destination.suffix.lower()
    writers: dict[str, Callable[[Path], None]] = {
        ".pdf": lambda target: _write_pdf(target, content),
        ".xlsx": lambda target: _write_workbook(target, content),
        ".docx": lambda target: _write_docx(target, content),
    }
    if extension in TEXT_EXT | {".csv", ".tsv"}:
        writer = lambda target: _write_text(target, content)
    else:
        writer = writers.get(extension)
    if writer is None:
        return (
            f"Unsupported output type {extension or '(none)'}. create_file supports "
            "PDF, XLSX, DOCX, CSV, TSV, and text-based files."
        )
    try:
        _atomic_create(destination, writer)
    except Exception as exc:
        return f"Could not create {destination}: {exc}"
    _FILES_READ.add(str(destination.resolve()))
    return f"Created {destination} ({destination.stat().st_size} bytes)."


def edit_file(
    path: str,
    operation: str,
    content: str,
    old_text: str = "",
    replace_all: bool = False,
) -> str:
    """Append to or exactly replace text in a file that the agent has read."""
    from .. import tools

    resolution = path_search.resolve_existing(path, "file")
    if resolution.path is None:
        return resolution.message
    target = resolution.path
    blocked = tools.write_path_block_reason(target)
    if blocked:
        return blocked
    if kind(target) != "text" and not _looks_like_text(target):
        return f"{target} is not an editable text file."
    key = str(target.resolve())
    if key not in _FILES_READ:
        return f"Read {target} with read_file before editing it."
    try:
        existing = target.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        return f"Could not read {target}: {exc}"

    normalized = str(operation or "").casefold().strip()
    if normalized == "append":
        separator = "" if not existing or existing.endswith("\n") or content.startswith("\n") else "\n"
        updated = existing + separator + content
        summary = f"Appended {len(content)} characters"
    elif normalized == "replace":
        if not old_text:
            return "old_text is required for replace."
        count = existing.count(old_text)
        if count == 0:
            return f"Text was not found in {target}; read it again and copy the exact text."
        if count > 1 and not replace_all:
            return (
                f"Found {count} matches in {target}. Provide more surrounding text "
                "or set replace_all=true."
            )
        updated = (
            existing.replace(old_text, content)
            if replace_all else existing.replace(old_text, content, 1)
        )
        summary = f"Replaced {count if replace_all else 1} occurrence(s)"
    else:
        return "operation must be append or replace."
    try:
        _atomic_create(target, lambda temporary: temporary.write_text(updated, encoding="utf-8"))
    except Exception as exc:
        return f"Could not edit {target}: {exc}"
    return f"{summary} in {target}."


def _operation_url(base_url: str, operation: str) -> str:
    parsed = urlsplit(str(base_url or "").strip())
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        raise RuntimeError("the generation model needs an HTTP(S) base URL")
    suffix = "/" + operation.strip("/")
    path = parsed.path.rstrip("/")
    if not path.endswith(suffix):
        path += suffix
    return urlunsplit((parsed.scheme, parsed.netloc, path, parsed.query, parsed.fragment))


def _provider_headers(runtime: dict) -> dict[str, str]:
    headers = {"Accept": "application/json", "Content-Type": "application/json"}
    if runtime.get("api_key"):
        headers["Authorization"] = f"Bearer {runtime['api_key']}"
    return headers


def _download_generated_image(url: str, headers: dict[str, str]) -> bytes:
    if url.startswith("data:") and ";base64," in url:
        return base64.b64decode(url.split(";base64,", 1)[1])
    download_headers = {
        key: value for key, value in headers.items() if key != "Content-Type"
    }
    download = requests.get(url, headers=download_headers, timeout=120)
    download.raise_for_status()
    return download.content


def _openai_image_bytes(runtime: dict, prompt: str) -> bytes:
    headers = _provider_headers(runtime)
    response = requests.post(
        _operation_url(runtime["base_url"], "images/generations"),
        headers=headers,
        json={"model": runtime["model"], "prompt": prompt, "response_format": "b64_json"},
        timeout=300,
    )
    if response.status_code >= 400:
        raise RuntimeError(
            f"generation endpoint returned HTTP {response.status_code}: "
            f"{response.text[:500]}"
        )
    body = response.json()
    item = (body.get("data") or [{}])[0]
    if item.get("b64_json"):
        return base64.b64decode(item["b64_json"])
    elif item.get("url"):
        return _download_generated_image(item["url"], headers)
    raise RuntimeError("generation endpoint returned neither b64_json nor a URL")


def _mistral_image_bytes(runtime: dict, prompt: str) -> bytes:
    """Use Mistral's provider-hosted image_generation built-in tool."""
    headers = _provider_headers(runtime)
    response = requests.post(
        _operation_url(runtime["base_url"], "chat/completions"),
        headers=headers,
        json={
            "model": runtime["model"],
            "messages": [{"role": "user", "content": prompt}],
            "tools": [{"type": "image_generation"}],
        },
        timeout=300,
    )
    if response.status_code >= 400:
        raise RuntimeError(
            f"Mistral image generation returned HTTP {response.status_code}: "
            f"{response.text[:500]}"
        )
    body = response.json()
    urls = []
    for choice in body.get("choices") or []:
        messages = choice.get("messages") or [choice.get("message") or {}]
        for message in messages:
            content = message.get("content") if isinstance(message, dict) else None
            chunks = content if isinstance(content, list) else []
            for chunk in chunks:
                if not isinstance(chunk, dict) or chunk.get("type") != "image_url":
                    continue
                image_url = chunk.get("image_url")
                if isinstance(image_url, dict):
                    image_url = image_url.get("url")
                if image_url:
                    urls.append(str(image_url))
            if isinstance(content, str):
                urls.extend(re.findall(r"https?://[^\s\])]+", content))
    if not urls:
        raise RuntimeError("Mistral returned no generated image URL")
    return _download_generated_image(urls[0], headers)


def _generate_image(destination: Path, prompt: str) -> None:
    from .. import db

    runtime = db.get_builtin_agent_generation_runtime("media")
    if runtime is None:
        raise RuntimeError(
            "select an image-generation model for Files and Media in Agent Studio"
        )
    provider = str(runtime.get("provider") or "").casefold()
    hostname = (urlsplit(runtime.get("base_url") or "").hostname or "").casefold()
    if "mistral" in provider or hostname.endswith("mistral.ai"):
        raw = _mistral_image_bytes(runtime, prompt)
    else:
        raw = _openai_image_bytes(runtime, prompt)
    try:
        from PIL import Image

        image = Image.open(io.BytesIO(raw))
        output_format = {
            ".jpg": "JPEG", ".jpeg": "JPEG", ".png": "PNG", ".webp": "WEBP",
            ".gif": "GIF", ".bmp": "BMP",
        }[destination.suffix.lower()]
        if output_format == "JPEG" and image.mode not in ("RGB", "L"):
            image = image.convert("RGB")
        image.save(destination, format=output_format)
    except ImportError:
        destination.write_bytes(raw)


def generate_media(path: str, prompt: str, specification: str = "") -> str:
    from .. import tools

    resolution = path_search.resolve_output(path)
    if resolution.path is None:
        return resolution.message
    destination = resolution.path
    blocked = tools.write_path_block_reason(destination)
    if blocked:
        return blocked
    extension = destination.suffix.lower()
    if extension in PRESENTATION_EXT:
        writer = lambda target: _write_presentation(target, specification or prompt)
    elif extension in IMAGE_EXT:
        writer = lambda target: _generate_image(target, prompt)
    elif extension in VIDEO_EXT:
        return (
            "Video generation is not available yet: compatible providers do not share "
            "one standard generation endpoint. Add a provider adapter before enabling it."
        )
    else:
        return (
            "generate_media supports images and PPTX presentations; video needs "
            "a provider adapter."
        )
    try:
        _atomic_create(destination, writer)
    except Exception as exc:
        return f"Could not generate {destination}: {exc}"
    return f"Generated {destination} ({destination.stat().st_size} bytes)."
